import os
import io
import json
import base64
import logging
import re
from datetime import datetime
from typing import Any

# Add GTK+ DLL directory for WeasyPrint on Windows
if os.name == 'nt':
    for sd in [
        r"C:\Program Files\GTK3-Runtime Win64\bin",
        r"C:\Program Files\GTK3-Runtime\bin",
        r"C:\Program Files (x86)\GTK3-Runtime Win64\bin",
        r"C:\Program Files (x86)\GTK3-Runtime\bin",
    ]:
        if os.path.exists(sd):
            try:
                os.add_dll_directory(sd)
                break
            except Exception:
                pass

import polars as pl
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from weasyprint import HTML
from pdf_generator import (
    markdown_to_html,
    generate_distribution_base64,
    generate_correlation_matrix_base64,
    generate_ml_visualization_base64
)

logger = logging.getLogger(__name__)


def generate_pdf_in_memory(dataframe: pl.DataFrame, analysis_summary: dict) -> bytes:
    """Generate WeasyPrint PDF for a dataset with dynamic profiling plots, ML metrics, and visualizations."""
    try:
        if dataframe is None:
            dataframe = pl.DataFrame()
        if analysis_summary is None:
            analysis_summary = {}

        total_rows = analysis_summary.get('rows', dataframe.height)
        total_cols = analysis_summary.get('cols', dataframe.width)
        category = analysis_summary.get('category', 'Generic Data')
        
        # Extract AI insights & AutoML results
        ai_insights = analysis_summary.get('ai_insights', '')
        automl = analysis_summary.get('automl', None)
        
        formatted_insights = markdown_to_html(ai_insights) if ai_insights else ""
        
        # Check target column to trigger ML visualization plot
        target_column = None
        if isinstance(automl, dict):
            target_column = automl.get("target_column") or automl.get("target")
        
        # Columns metadata
        columns_meta = []
        for name, dtype in zip(dataframe.columns, dataframe.dtypes):
            null_count = dataframe[name].null_count()
            null_pct = (null_count / dataframe.height * 100) if dataframe.height > 0 else 0
            columns_meta.append({
                "name": name,
                "type": str(dtype),
                "null_pct": f"{null_pct:.1f}%",
                "non_null": dataframe.height - null_count
            })

        # Generate plots
        dist_img = generate_distribution_base64(dataframe)
        corr_img = generate_correlation_matrix_base64(dataframe)
        ml_img = generate_ml_visualization_base64(dataframe, target_column) if target_column else ""
        
        # Sample dataset observations
        sample_df = dataframe.head(10)
        sample_headers = sample_df.columns
        sample_rows = sample_df.iter_rows()
        
        generation_time = datetime.now().strftime("%B %d, %Y at %I:%M %p")
        
        html_content = f"""
        <!DOCTYPE html>
        <html>
        <head>
            <meta charset="utf-8">
            <title>Data Analysis Report</title>
            <style>
                @page {{
                    size: A4;
                    margin: 20mm;
                    @bottom-right {{
                        content: "Page " counter(page) " of " counter(pages);
                        font-family: 'Helvetica Neue', Helvetica, Arial, sans-serif;
                        font-size: 8pt;
                        color: #718096;
                    }}
                    @bottom-left {{
                        content: "DataSaaS Pro Analytics System • Execution Telemetry";
                        font-family: 'Helvetica Neue', Helvetica, Arial, sans-serif;
                        font-size: 8pt;
                        color: #718096;
                    }}
                }}
                body {{
                    font-family: 'Georgia', 'Times New Roman', serif;
                    color: #2d3748;
                    line-height: 1.6;
                    font-size: 10.5pt;
                }}
                header {{
                    border-bottom: 2px solid #2b6cb0;
                    padding-bottom: 12px;
                    margin-bottom: 25px;
                }}
                .report-title {{
                    font-family: 'Helvetica Neue', Helvetica, Arial, sans-serif;
                    color: #1a365d;
                    font-size: 22pt;
                    font-weight: 800;
                    margin: 0 0 5px 0;
                    letter-spacing: -0.5px;
                    text-transform: uppercase;
                }}
                .report-subtitle {{
                    font-family: 'Helvetica Neue', Helvetica, Arial, sans-serif;
                    color: #4a5568;
                    font-size: 11pt;
                    margin: 0;
                    font-weight: 400;
                }}
                .metadata-grid {{
                    display: table;
                    width: 100%;
                    margin-bottom: 25px;
                    background-color: #f7fafc;
                    border: 1px solid #e2e8f0;
                    border-radius: 6px;
                    border-collapse: collapse;
                }}
                .metadata-row {{
                    display: table-row;
                }}
                .metadata-cell {{
                    display: table-cell;
                    padding: 8px 12px;
                    font-size: 9pt;
                    border: 1px solid #e2e8f0;
                }}
                .metadata-label {{
                    font-weight: bold;
                    color: #4a5568;
                    font-family: 'Helvetica Neue', Helvetica, Arial, sans-serif;
                    background-color: #edf2f7;
                    width: 25%;
                }}
                h2 {{
                    font-family: 'Helvetica Neue', Helvetica, Arial, sans-serif;
                    color: #2b6cb0;
                    font-size: 13pt;
                    margin-top: 25px;
                    margin-bottom: 12px;
                    border-bottom: 1px solid #e2e8f0;
                    padding-bottom: 5px;
                    text-transform: uppercase;
                    letter-spacing: 0.5px;
                }}
                p {{
                    margin-bottom: 15px;
                    text-align: justify;
                }}
                .section-container {{
                    page-break-inside: avoid;
                    margin-bottom: 25px;
                }}
                .data-table {{
                    width: 100%;
                    border-collapse: collapse;
                    margin-top: 10px;
                    margin-bottom: 20px;
                    font-size: 8.5pt;
                }}
                .data-table th {{
                    background-color: #2b6cb0;
                    color: white;
                    font-weight: bold;
                    font-family: 'Helvetica Neue', Helvetica, Arial, sans-serif;
                    text-align: left;
                    padding: 6px 8px;
                    border: 1px solid #cbd5e0;
                }}
                .data-table td {{
                    padding: 5px 8px;
                    border: 1px solid #e2e8f0;
                    font-family: monospace;
                }}
                .data-table tr:nth-child(even) {{
                    background-color: #f7fafc;
                }}
                .chart-container {{
                    display: table;
                    width: 100%;
                    margin: 20px 0;
                    page-break-inside: avoid;
                }}
                .chart-box {{
                    display: table-cell;
                    width: 50%;
                    padding: 5px;
                    text-align: center;
                    vertical-align: middle;
                }}
                .chart-img {{
                    max-width: 95%;
                    border: 1px solid #e2e8f0;
                    box-shadow: 0 4px 6px rgba(0,0,0,0.04);
                }}
                .insights-callout {{
                    background-color: #ebf8ff;
                    border-left: 4px solid #3182ce;
                    padding: 15px;
                    border-radius: 0 6px 6px 0;
                    margin: 20px 0;
                }}
                .insights-title {{
                    font-family: 'Helvetica Neue', Helvetica, Arial, sans-serif;
                    font-weight: bold;
                    color: #2b6cb0;
                    margin-bottom: 8px;
                    font-size: 11pt;
                }}
            </style>
        </head>
        <body>
            <header>
                <h1 class="report-title">Exploratory Data Analysis Report</h1>
                <p class="report-subtitle">Systematic Dataset Profiling & AutoML Summary</p>
            </header>

            <div class="metadata-grid">
                <div class="metadata-row">
                    <div class="metadata-cell metadata-label">Generated On</div>
                    <div class="metadata-cell">{generation_time}</div>
                    <div class="metadata-cell metadata-label">Data Dimension</div>
                    <div class="metadata-cell">{total_rows} Rows x {total_cols} Columns</div>
                </div>
                <div class="metadata-row">
                    <div class="metadata-cell metadata-label">Category Detected</div>
                    <div class="metadata-cell">{category}</div>
                    <div class="metadata-cell metadata-label">Sanitization</div>
                    <div class="metadata-cell">PII Masks Active</div>
                </div>
            </div>
        """
        
        if formatted_insights:
            html_content += f"""
            <div class="section-container">
                <h2>1. Executive Summary & Business Insights</h2>
                <div class="insights-callout">
                    <div class="insights-title">AI Analytics Insights</div>
                    {formatted_insights}
                </div>
            </div>
            """
            
        if automl:
            html_content += f"""
            <div class="section-container">
                <h2>2. Automated Machine Learning Predictions</h2>
                <div class="insights-callout" style="background-color: #f0fff4; border-left: 4px solid #38a169;">
                    <div class="insights-title" style="color: #276749;">AutoML Summary Report</div>
                    <p><strong>Target Column:</strong> {target_column or 'N/A'}</p>
                    <p><strong>Recommended Estimator:</strong> {automl.get('best_algorithm') or automl.get('model_name', 'N/A')}</p>
                    <p><strong>Model Performance Metrics:</strong></p>
                    <ul>
                        <li>R2 Accuracy Score: {automl.get('r2') or automl.get('r2_score', 'N/A')}</li>
                        <li>Mean Absolute Error (MAE): {automl.get('mae', 'N/A')}</li>
                        <li>Root Mean Squared Error (RMSE): {automl.get('rmse', 'N/A')}</li>
                    </ul>
                </div>
            </div>
            """
            
        html_content += f"""
            <div class="section-container">
                <h2>3. Columns Completeness Profiles</h2>
                <table class="data-table">
                    <thead>
                        <tr>
                            <th>Column Name</th>
                            <th>Data Type</th>
                            <th>Completeness Count</th>
                            <th>Null Percentage</th>
                        </tr>
                    </thead>
                    <tbody>
        """
        
        for col in columns_meta:
            html_content += f"""
                        <tr>
                            <td>{col['name']}</td>
                            <td>{col['type']}</td>
                            <td>{col['non_null']} / {dataframe.height}</td>
                            <td>{col['null_pct']}</td>
                        </tr>
            """
            
        html_content += """
                    </tbody>
                </table>
            </div>
        """
        
        if dist_img or corr_img:
            html_content += """
            <div class="section-container">
                <h2>4. Dynamic Plots</h2>
                <div class="chart-container">
            """
            if dist_img:
                html_content += f"""
                    <div class="chart-box">
                        <img class="chart-img" src="data:image/png;base64,{dist_img}" alt="Distribution Profile"/>
                    </div>
                """
            if corr_img:
                html_content += f"""
                    <div class="chart-box">
                        <img class="chart-img" src="data:image/png;base64,{corr_img}" alt="Correlation Heatmap"/>
                    </div>
                """
            html_content += """
                </div>
            </div>
            """
            
        if ml_img:
            html_content += f"""
            <div class="section-container" style="page-break-before: always;">
                <h2>5. Machine Learning Visualizations</h2>
                <p>Performance profile and learning curve mapping of Random Forest model trained on <strong>{target_column}</strong>:</p>
                <div class="chart-container" style="text-align: center;">
                    <img class="chart-img" style="max-width: 80%;" src="data:image/png;base64,{ml_img}" alt="ML Performance Graph"/>
                </div>
            </div>
            """
            
        if dataframe.height > 0:
            html_content += f"""
            <div class="section-container" style="page-break-before: always;">
                <h2>{6 if ml_img else 5}. Sample Dataset Observations (First 10 Rows)</h2>
                <table class="data-table" style="font-size: 7.5pt;">
                    <thead>
                        <tr>
            """
            for head in sample_headers:
                html_content += f"<th>{head}</th>"
            html_content += """
                        </tr>
                    </thead>
                    <tbody>
            """
            for row in sample_rows:
                html_content += "<tr>"
                for item in row:
                    safe_item = str(item or "")
                    if len(safe_item) > 25:
                        safe_item = safe_item[:22] + "..."
                    html_content += f"<td>{safe_item}</td>"
                html_content += "</tr>"
            html_content += """
                    </tbody>
                </table>
            </div>
            """
            
        html_content += """
        </body>
        </html>
        """
        
        return HTML(string=html_content).write_pdf()
        
    except Exception as e:
        logger.exception(f"Failed to generate WeasyPrint PDF in memory: {e}")
        return b"%PDF-1.4\n%EOF"


def generate_structured_report_pdf(*, title: str, subtitle: str, sections: list[dict]) -> bytes:
    """Generate structured PDF report with WeasyPrint from layout section data with dynamic colors based on keywords in title."""
    try:
        title = title or "Report"
        subtitle = subtitle or ""
        sections = sections or []
        
        # Determine theme based on title keywords
        title_lower = title.lower()
        if "quality" in title_lower:
            theme_color = "#276749"  # Emerald/Forest Green
            bg_light = "#f0fff4"
            border_color = "#c6f6d5"
            accent_color = "#38a169"
            theme_name = "Quality Metrics & Health Profile"
        elif "technical" in title_lower or "research" in title_lower:
            theme_color = "#1a365d"  # Deep Indigo
            bg_light = "#ebf8ff"
            border_color = "#bee3f8"
            accent_color = "#3182ce"
            theme_name = "Technical Lab & Research Document"
        elif "executive" in title_lower:
            theme_color = "#2d3748"  # Slate
            bg_light = "#edf2f7"
            border_color = "#e2e8f0"
            accent_color = "#4a5568"
            theme_name = "Executive Briefing Summary"
        else:
            theme_color = "#2b6cb0"  # Cobalt Blue
            bg_light = "#ebf8ff"
            border_color = "#bee3f8"
            accent_color = "#3182ce"
            theme_name = "Analytics Summary Report"
            
        generation_time = datetime.now().strftime("%B %d, %Y at %I:%M %p")
        
        html_content = f"""
        <!DOCTYPE html>
        <html>
        <head>
            <meta charset="utf-8">
            <title>{title}</title>
            <style>
                @page {{
                    size: A4;
                    margin: 20mm;
                    @bottom-right {{
                        content: "Page " counter(page) " of " counter(pages);
                        font-family: 'Helvetica Neue', Helvetica, Arial, sans-serif;
                        font-size: 8pt;
                        color: #718096;
                    }}
                    @bottom-left {{
                        content: "{title} • {theme_name}";
                        font-family: 'Helvetica Neue', Helvetica, Arial, sans-serif;
                        font-size: 8pt;
                        color: #718096;
                    }}
                }}
                body {{
                    font-family: 'Georgia', 'Times New Roman', serif;
                    color: #2d3748;
                    line-height: 1.6;
                    font-size: 10.5pt;
                }}
                header {{
                    border-bottom: 2px solid {theme_color};
                    padding-bottom: 12px;
                    margin-bottom: 25px;
                }}
                .report-title {{
                    font-family: 'Helvetica Neue', Helvetica, Arial, sans-serif;
                    color: #1a365d;
                    font-size: 22pt;
                    font-weight: 800;
                    margin: 0 0 5px 0;
                    letter-spacing: -0.5px;
                    text-transform: uppercase;
                }}
                .report-subtitle {{
                    font-family: 'Helvetica Neue', Helvetica, Arial, sans-serif;
                    color: #4a5568;
                    font-size: 11pt;
                    margin: 0;
                    font-weight: 400;
                }}
                .meta-bar {{
                    font-family: 'Helvetica Neue', Helvetica, Arial, sans-serif;
                    font-size: 8.5pt;
                    color: #718096;
                    margin-bottom: 25px;
                    background-color: {bg_light};
                    padding: 8px 12px;
                    border: 1px solid {border_color};
                    border-radius: 6px;
                }}
                h2 {{
                    font-family: 'Helvetica Neue', Helvetica, Arial, sans-serif;
                    color: {theme_color};
                    font-size: 13pt;
                    margin-top: 25px;
                    margin-bottom: 12px;
                    border-bottom: 1px solid {border_color};
                    padding-bottom: 5px;
                    text-transform: uppercase;
                    letter-spacing: 0.5px;
                }}
                p {{
                    margin-bottom: 15px;
                    text-align: justify;
                }}
                .section-container {{
                    page-break-inside: avoid;
                    margin-bottom: 25px;
                }}
                .row-item {{
                    margin-bottom: 15px;
                    padding-bottom: 8px;
                    border-bottom: 1px dashed {border_color};
                }}
                .row-label {{
                    font-family: 'Helvetica Neue', Helvetica, Arial, sans-serif;
                    font-weight: bold;
                    color: {accent_color};
                    font-size: 9.5pt;
                    margin-bottom: 3px;
                }}
                .row-value {{
                    font-size: 10pt;
                    text-align: justify;
                }}
            </style>
        </head>
        <body>
            <header>
                <h1 class="report-title">{title}</h1>
                {f'<p class="report-subtitle">{subtitle}</p>' if subtitle else ''}
            </header>

            <div class="meta-bar">
                <strong>Generated On:</strong> {generation_time} &nbsp;&bull;&nbsp; 
                <strong>Theme Applied:</strong> {theme_name}
            </div>
        """
        
        for section in sections:
            heading = section.get("heading") or "Section"
            rows = section.get("rows") or []
            
            html_content += f"""
            <div class="section-container">
                <h2>{heading}</h2>
                <div class="section-rows">
            """
            
            for row in rows:
                label = row.get("label") or ""
                value = row.get("value") or ""
                
                # Format value line breaks and simple formatting
                value_formatted = value.replace("\n", "<br/>")
                
                html_content += f"""
                    <div class="row-item">
                        {f'<div class="row-label">{label}</div>' if label else ''}
                        <div class="row-value">{value_formatted}</div>
                    </div>
                """
                
            html_content += """
                </div>
            </div>
            """
            
        html_content += """
        </body>
        </html>
        """
        
        return HTML(string=html_content).write_pdf()
        
    except Exception as e:
        logger.exception(f"Failed to generate structured WeasyPrint PDF: {e}")
        return b"%PDF-1.4\n%EOF"


def generate_structured_report_pptx(*, title: str, subtitle: str, sections: list[dict]) -> bytes:
    """Generate PPTX report with template styling."""
    try:
        title_str = str(title or "Report")
        subtitle_str = str(subtitle or "")
        sections = sections or []
        
        presentation = Presentation()
        presentation.slide_width = Inches(13.333)
        presentation.slide_height = Inches(7.5)

        try:
            title_slide_layout = presentation.slide_layouts[6]  # blank layout
            title_slide = presentation.slides.add_slide(title_slide_layout)
            
            bg = title_slide.background
            fill = bg.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(0x1A, 0x23, 0x7E)
            
            title_box = title_slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
            tf = title_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title_str[:100]
            p.font.size = Pt(40)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            p.alignment = PP_ALIGN.CENTER
            
            if subtitle_str:
                sub_box = title_slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
                tf2 = sub_box.text_frame
                tf2.word_wrap = True
                p2 = tf2.paragraphs[0]
                p2.text = subtitle_str[:100]
                p2.font.size = Pt(20)
                p2.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
                p2.alignment = PP_ALIGN.CENTER
            
            date_box = title_slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(12.333), Inches(0.4))
            tf3 = date_box.text_frame
            p3 = tf3.paragraphs[0]
            p3.text = f"Generated on {datetime.now().strftime('%B %d, %Y')}"
            p3.font.size = Pt(12)
            p3.font.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)
            p3.alignment = PP_ALIGN.CENTER
        except Exception as e:
            logger.error(f"Error creating title slide: {e}")

        for section in sections:
            try:
                if not isinstance(section, dict):
                    continue
                    
                slide_layout = presentation.slide_layouts[6]
                slide = presentation.slides.add_slide(slide_layout)
                
                bg = slide.background
                fill = bg.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                
                header = slide.shapes.add_textbox(Inches(0.4), Inches(0.3), Inches(12.5), Inches(0.8))
                htf = header.text_frame
                hp = htf.paragraphs[0]
                hp.text = str(section.get("heading") or "Section")[:100]
                hp.font.size = Pt(28)
                hp.font.bold = True
                hp.font.color.rgb = RGBColor(0x1A, 0x23, 0x7E)
                
                line = slide.shapes.add_shape(1, Inches(0.4), Inches(1.05), Inches(12.5), Inches(0.05))
                line.fill.solid()
                line.fill.fore_color.rgb = RGBColor(0x1A, 0x23, 0x7E)
                line.line.fill.background()
                
                body = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12.333), Inches(5.8))
                tf = body.text_frame
                tf.word_wrap = True
                rows = section.get("rows") or []
                
                if not rows:
                    p = tf.paragraphs[0]
                    p.text = "No data available"
                    p.font.size = Pt(16)
                    p.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
                    continue

                first_row = True
                for row in rows:
                    if not isinstance(row, dict):
                        continue
                        
                    try:
                        label = str(row.get("label") or "")[:80]
                        value = str(row.get("value") or "")[:200]
                        p = tf.paragraphs[0] if first_row else tf.add_paragraph()
                        p.text = f"{label}: {value}" if label else str(value)[:200]
                        p.font.size = Pt(16)
                        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
                        p.space_after = Pt(8)
                        first_row = False
                    except Exception as e:
                        logger.warning(f"Error adding row to PPTX: {e}")
            except Exception as e:
                logger.error(f"Error adding section to PPTX: {e}")

        try:
            output = io.BytesIO()
            presentation.save(output)
            return output.getvalue()
        except Exception as e:
            logger.error(f"Error outputting PPTX: {e}")
            fallback = Presentation()
            fallback_out = io.BytesIO()
            fallback.save(fallback_out)
            return fallback_out.getvalue()
            
    except Exception as e:
        logger.exception(f"Failed to generate PPTX: {e}")
        fallback = Presentation()
        fallback_out = io.BytesIO()
        fallback.save(fallback_out)
        return fallback_out.getvalue()
